import os
import json
import uuid
from pathlib import Path
from pptx import Presentation
import yake

# ==== 설정 ====
PPT_PATH = "./data/manual_docs/psi_manual.pptx"
CHUNK_SAVE_PATH = "./data/rag_chunks/psi_slide_chunks.json"
LANGUAGE = "ko"  # 한국어
MAX_KEYWORDS = 5

# 키워드 추출기 초기화
kw_extractor = yake.KeywordExtractor(
    lan=LANGUAGE,
    n=1,
    top=MAX_KEYWORDS,
    stopwords=None  # 영어 stopword 제외 안 함
)

# 키워드 추출 함수
def extract_keywords(text):
    keywords = kw_extractor.extract_keywords(text)
    return [kw for kw, _ in keywords]

# 슬라이드 텍스트 → 청크로 분할 및 키워드 추출
def extract_text_chunks(ppt_path):
    prs = Presentation(ppt_path)
    chunks = []

    for idx, slide in enumerate(prs.slides, start=1):
        full_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"

        # 줄 단위 청크
        for para in full_text.strip().split("\n"):
            clean_text = para.strip()
            if clean_text:
                keywords = extract_keywords(clean_text)
                chunks.append({
                    "id": str(uuid.uuid4()),
                    "slide_index": idx,
                    "text": clean_text,
                    "source": f"slide_{idx}",
                    "keywords": keywords
                })

    return chunks

# ==== 실행 ====
if __name__ == "__main__":
    print("🔍 슬라이드 텍스트 기반 RAG 청크 생성 중...")

    chunks = extract_text_chunks(PPT_PATH)

    os.makedirs(os.path.dirname(CHUNK_SAVE_PATH), exist_ok=True)
    with open(CHUNK_SAVE_PATH, "w", encoding="utf-8") as f:
        json.dump(chunks, f, ensure_ascii=False, indent=2)

    print(f"✅ 총 {len(chunks)}개의 청크가 저장되었습니다.")
    print(f"📁 저장 위치: {CHUNK_SAVE_PATH}")
