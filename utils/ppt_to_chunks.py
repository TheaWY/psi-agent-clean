import os
import json
import uuid
from pathlib import Path
from pptx import Presentation
from win32com.client import Dispatch

# ===== 경로 설정 =====
PPT_PATH = "C:/Users/LGCNS/Desktop/code/psi_agent_system_gpt_only_final_complete/data/manual_docs/psi_manual.pptx"
CHUNK_JSON_PATH = "C:/Users/LGCNS/Desktop/code/psi_agent_system_gpt_only_final_complete/data/rag_chunks/psi_slide_chunks.json"
IMAGE_DIR = Path("C:/Users/LGCNS/Desktop/code/psi_agent_system_gpt_only_final_complete/data/rag_chunks/images")
IMAGE_DIR.mkdir(parents=True, exist_ok=True)

# ===== 슬라이드 이미지 저장 =====
def export_slide_images(ppt_path, output_dir):
    powerpoint = Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
    presentation.SaveAs(str(output_dir), 17)  # 17 = pptExportAsPNG
    presentation.Close()
    powerpoint.Quit()

    # Slide1.PNG -> slide_1.png 형식으로 리네이밍
    for img_file in Path(output_dir).glob("Slide*.PNG"):
        slide_num = img_file.stem.replace("Slide", "").lstrip("0") or "0"
        new_name = f"slide_{slide_num}.png"
        img_file.rename(img_file.with_name(new_name))

# ===== 슬라이드 텍스트 추출 및 청크 분할 =====
def extract_slide_text_chunks(pptx_path, max_chunk_length=300):
    prs = Presentation(pptx_path)
    chunks = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        full_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)

        slide_text = "\n".join(full_text).strip()
        slide_lines = [line.strip() for line in slide_text.splitlines() if line.strip()]

        chunk_buffer = ""
        chunk_index = 1
        for line in slide_lines:
            if len(chunk_buffer) + len(line) < max_chunk_length:
                chunk_buffer += line + "\n"
            else:
                chunks.append({
                    "id": str(uuid.uuid4()),
                    "slide_index": slide_index,
                    "chunk_index": chunk_index,
                    "text": chunk_buffer.strip(),
                    "source": f"slide_{slide_index}",
                    "image_path": str(IMAGE_DIR / f"slide_{slide_index}.png").replace("\\", "/")
                })
                chunk_index += 1
                chunk_buffer = line + "\n"

        if chunk_buffer:
            chunks.append({
                "id": str(uuid.uuid4()),
                "slide_index": slide_index,
                "chunk_index": chunk_index,
                "text": chunk_buffer.strip(),
                "source": f"slide_{slide_index}",
                "image_path": str(IMAGE_DIR / f"slide_{slide_index}.png").replace("\\", "/")
            })

    return chunks

# ===== 실행 흐름 =====
print("🔍 PPT 슬라이드 이미지 및 텍스트 추출 중...")
export_slide_images(PPT_PATH, IMAGE_DIR)
chunks = extract_slide_text_chunks(PPT_PATH)

print("💾 JSON 저장 중...")
with open(CHUNK_JSON_PATH, "w", encoding="utf-8") as f:
    json.dump(chunks, f, ensure_ascii=False, indent=2)

print(f"✅ 완료: {len(chunks)}개 청크 저장됨 → {CHUNK_JSON_PATH}")